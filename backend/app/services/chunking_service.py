"""
Capa de chunking (división de texto) de la Capa de Trabajo de IA.

Herramienta utilizada: langchain-text-splitters (RecursiveCharacterTextSplitter).
Es el splitter más utilizado en arquitecturas RAG: intenta cortar primero por
párrafos, luego por líneas, luego por oraciones y, como último recurso, por
caracteres. Esto preserva mejor el sentido de cada fragmento que un corte de
tamaño fijo.

Extracción de texto por tipo de documento:
- PDF  -> pypdf (extracción nativa por página, sin OCR).
- PPTX -> python-pptx (extracción por diapositiva, usa el título si existe).
- TXT  -> lectura directa del archivo completo.
"""

import os
from pathlib import Path

from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pypdf import PdfReader

CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "900"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "150"))

_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    separators=["\n\n", "\n", ". ", " ", ""],
)


def extract_pdf_pages(file_path: Path) -> list[dict]:
    reader = PdfReader(str(file_path))
    pages = []

    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""

        pages.append({
            "page_number": index,
            "section_title": None,
            "text": text,
        })

    return pages


def extract_pptx_slides(file_path: Path) -> list[dict]:
    presentation = Presentation(str(file_path))
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        title = None

        title_shape = slide.shapes.title if slide.shapes else None

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            shape_text = "\n".join(
                paragraph.text
                for paragraph in shape.text_frame.paragraphs
                if paragraph.text
            ).strip()

            if not shape_text:
                continue

            if title is None and title_shape is not None and shape == title_shape:
                title = shape_text

            texts.append(shape_text)

        slides.append({
            "page_number": index,
            "section_title": title,
            "text": "\n".join(texts),
        })

    return slides


def extract_txt_pages(file_path: Path) -> list[dict]:
    text = file_path.read_text(encoding="utf-8", errors="ignore")

    return [{
        "page_number": None,
        "section_title": None,
        "text": text,
    }]


def extract_pages(file_path: Path, mime_type: str | None) -> list[dict]:
    suffix = file_path.suffix.lower()

    if suffix == ".pdf" or (mime_type and "pdf" in mime_type):
        return extract_pdf_pages(file_path)

    if suffix == ".pptx" or (mime_type and "presentation" in mime_type):
        return extract_pptx_slides(file_path)

    if suffix == ".txt" or (mime_type and "text/plain" in mime_type):
        return extract_txt_pages(file_path)

    raise ValueError(
        f"Tipo de archivo no soportado para chunking: {file_path.name}"
    )


def chunk_document(file_path: Path, mime_type: str | None) -> list[dict]:
    """Extrae el texto del documento y lo divide en fragmentos (chunks).

    Devuelve una lista de dicts con: chunk_index, content, page_number y
    section_title, listos para insertarse en la tabla document_chunks.
    """

    pages = extract_pages(file_path, mime_type)

    chunks = []
    chunk_index = 0

    for page in pages:
        page_text = (page.get("text") or "").strip()

        if not page_text:
            continue

        for fragment in _splitter.split_text(page_text):
            clean_fragment = fragment.strip()

            if not clean_fragment:
                continue

            chunks.append({
                "chunk_index": chunk_index,
                "content": clean_fragment,
                "page_number": page.get("page_number"),
                "section_title": page.get("section_title"),
            })

            chunk_index += 1

    return chunks
